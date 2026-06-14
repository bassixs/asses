from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from bot.services.reports import LEVEL_DESCRIPTIONS, _format_course, _nearest_level_description

_ACCENT = RGBColor(0x1F, 0x4E, 0x79)
_DARK = RGBColor(0x22, 0x22, 0x22)
_GREY = RGBColor(0x55, 0x55, 0x55)
_BODY_TOP = Inches(1.25)
_BODY_HEIGHT = Inches(5.9)
_BODY_LEFT = Inches(0.6)
_BODY_WIDTH = Inches(12.1)
# Rough budgets so long content paginates across slides instead of overflowing.
_CHAR_BUDGET = 1500
_MAX_BULLETS = 12


def save_participant_report_pptx(
    *,
    path: Path,
    participant_name: str,
    center_name: str | None,
    exercise_names: list[str],
    report_json: dict[str, Any],
) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    competencies: dict[str, Any] = report_json.get("competencies", {}) or {}
    matrix: dict[str, Any] = report_json.get("matrix", {}) or {}
    date_str = str(report_json.get("generated_date") or "")

    _title_slide(prs, participant_name, center_name, date_str)
    _about_slide(prs, exercise_names)
    _matrix_slide(prs, competencies, matrix, exercise_names)
    _summary_slide(prs, competencies)

    for competence, data in competencies.items():
        _competence_result_slide(prs, str(competence), data)
        _competence_recommendations_slide(prs, str(competence), data)
        _competence_practice_slide(prs, str(competence), data)

    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(path)


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _title_bar(slide, text: str) -> None:
    box = slide.shapes.add_textbox(_BODY_LEFT, Inches(0.4), _BODY_WIDTH, Inches(0.8))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.bold = True
    run.font.size = Pt(26)
    run.font.color.rgb = _ACCENT


def _body(slide):
    box = slide.shapes.add_textbox(_BODY_LEFT, _BODY_TOP, _BODY_WIDTH, _BODY_HEIGHT)
    tf = box.text_frame
    tf.word_wrap = True
    return tf


def _para(tf, text: str, *, size: int = 12, bold: bool = False, color=_DARK, bullet: bool = False, first: bool = False):
    paragraph = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    run = paragraph.add_run()
    run.text = (f"•  {text}" if bullet else text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    paragraph.space_after = Pt(4)
    return paragraph


def _title_slide(prs, participant_name: str, center_name: str | None, date_str: str) -> None:
    slide = _blank(prs)
    box = slide.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(11.3), Inches(2.5))
    tf = box.text_frame
    tf.word_wrap = True
    title = tf.paragraphs[0]
    title.alignment = PP_ALIGN.CENTER
    run = title.add_run()
    run.text = "ОТЧЕТ ПО РЕЗУЛЬТАТАМ ОЦЕНКИ КОМПЕТЕНЦИЙ"
    run.font.bold = True
    run.font.size = Pt(30)
    run.font.color.rgb = _ACCENT

    sub = tf.add_paragraph()
    sub.alignment = PP_ALIGN.CENTER
    r = sub.add_run()
    r.text = f"\nУчастник: {participant_name}"
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = _DARK
    if center_name:
        c = tf.add_paragraph()
        c.alignment = PP_ALIGN.CENTER
        rc = c.add_run()
        rc.text = f"Ассессмент-центр: {center_name}"
        rc.font.size = Pt(14)
        rc.font.color.rgb = _GREY
    if date_str:
        d = tf.add_paragraph()
        d.alignment = PP_ALIGN.CENTER
        rd = d.add_run()
        rd.text = f"Дата: {date_str}"
        rd.font.size = Pt(14)
        rd.font.color.rgb = _GREY

    note = slide.shapes.add_textbox(Inches(1.0), Inches(6.4), Inches(11.3), Inches(0.8))
    ntf = note.text_frame
    ntf.word_wrap = True
    np = ntf.paragraphs[0]
    np.alignment = PP_ALIGN.CENTER
    nr = np.add_run()
    nr.text = "Отчет конфиденциален. Не рекомендуется использовать его по истечении двух лет."
    nr.font.size = Pt(10)
    nr.font.italic = True
    nr.font.color.rgb = _GREY


def _about_slide(prs, exercise_names: list[str]) -> None:
    slide = _blank(prs)
    _title_bar(slide, "О центре оценки и развития")
    tf = _body(slide)
    _para(
        tf,
        "Центр оценки и развития — серия процедур для оценки компетенций руководителей в "
        "упражнениях, моделирующих ключевые моменты профессиональной деятельности.",
        size=13,
        first=True,
    )
    _para(tf, "Оценочные процедуры:", size=14, bold=True, color=_ACCENT)
    if exercise_names:
        for name in exercise_names:
            _para(tf, name, bullet=True)
    else:
        _para(tf, "Данные по упражнениям не указаны.")


def _matrix_slide(prs, competencies: dict[str, Any], matrix: dict[str, Any], exercise_names: list[str]) -> None:
    slide = _blank(prs)
    _title_bar(slide, "Матрица оценки и шкала")

    comps = list(competencies.keys())
    cols = ["Компетенция", *exercise_names, "Средний"]
    rows = len(comps) + 1
    table_height = Inches(min(4.5, 0.4 + 0.35 * rows))
    table = slide.shapes.add_table(rows, len(cols), _BODY_LEFT, _BODY_TOP, _BODY_WIDTH, table_height).table

    for ci, head in enumerate(cols):
        cell = table.cell(0, ci)
        cell.text = head
        _style_cell(cell, bold=True, size=11, color=RGBColor(0xFF, 0xFF, 0xFF), fill=_ACCENT)

    for ri, comp in enumerate(comps, start=1):
        _style_cell(table.cell(ri, 0), text=str(comp), bold=True, size=10)
        per_ex = matrix.get(comp, {}) or {}
        for ci, ex in enumerate(exercise_names, start=1):
            _style_cell(table.cell(ri, ci), text=_fmt_level(per_ex.get(ex)), size=10, align=PP_ALIGN.CENTER)
        avg = competencies.get(comp, {}).get("avg_level")
        _style_cell(table.cell(ri, len(cols) - 1), text=_fmt_level(avg), bold=True, size=10, align=PP_ALIGN.CENTER)

    note = slide.shapes.add_textbox(_BODY_LEFT, Inches(5.9), _BODY_WIDTH, Inches(1.2))
    ntf = note.text_frame
    ntf.word_wrap = True
    _para(ntf, "Требуемый уровень развития компетенции для успешного выполнения работы — 2 балла.", size=12, bold=True, first=True)
    _para(ntf, "Шкала 0–3 с половинными баллами; уровень рассчитывается по доле проявленных индикаторов (без «НЗ»).", size=11, color=_GREY)


def _summary_slide(prs, competencies: dict[str, Any]) -> None:
    slide = _blank(prs)
    _title_bar(slide, "Результаты центра оценки и развития")
    tf = _body(slide)
    first = True
    for comp, data in competencies.items():
        level = _fmt_level(data.get("avg_level"))
        p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
        first = False
        r1 = p.add_run()
        r1.text = f"{comp}:  "
        r1.font.size = Pt(14)
        r1.font.color.rgb = _DARK
        r2 = p.add_run()
        r2.text = level
        r2.font.size = Pt(14)
        r2.font.bold = True
        r2.font.color.rgb = _ACCENT
        p.space_after = Pt(6)


def _competence_result_slide(prs, competence: str, data: dict[str, Any]) -> None:
    slide = _blank(prs)
    _title_bar(slide, competence)
    tf = _body(slide)
    level = float(data.get("avg_level", 0) or 0)
    _para(tf, f"Ваш результат: {_fmt_level(level)}", size=16, bold=True, color=_ACCENT, first=True)
    _para(tf, _nearest_level_description(level), size=11, color=_GREY)
    _para(tf, "Сильные стороны", size=14, bold=True, color=_ACCENT)
    strengths = data.get("report_strengths") or data.get("strengths") or []
    if strengths:
        for item in strengths:
            _para(tf, item, bullet=True)
    else:
        _para(tf, "Не выявлено по обработанным упражнениям.", color=_GREY)
    _para(tf, "Зоны роста", size=14, bold=True, color=_ACCENT)
    zones = data.get("report_growth_zones") or data.get("growth_zones") or []
    if zones:
        for item in zones:
            _para(tf, item, bullet=True)
    else:
        _para(tf, "Не выявлено по обработанным упражнениям.", color=_GREY)


def _competence_recommendations_slide(prs, competence: str, data: dict[str, Any]) -> None:
    items: list[tuple[str, list[str]]] = []
    recommendations = data.get("recommendations") or []
    if recommendations:
        items.append(("Рекомендации по развитию", list(recommendations)))
    literature = data.get("literature") or []
    if literature:
        items.append(("Литература", list(literature)))
    courses = data.get("courses") or []
    if courses:
        items.append(("Тренинговые блоки", [_format_course(c) for c in courses]))
    _add_sectioned_slides(prs, f"{competence} — рекомендации", items)


def _competence_practice_slide(prs, competence: str, data: dict[str, Any]) -> None:
    tasks = data.get("practice_tasks") or []
    if not tasks:
        return
    _add_sectioned_slides(prs, f"{competence} — задания для закрепления", [("Задания для закрепления", list(tasks))])


def _add_sectioned_slides(prs, title: str, sections: list[tuple[str, list[str]]]) -> None:
    """Render labelled bullet sections, paginating onto extra slides when content is long."""
    flat: list[tuple[str | None, str]] = []
    for heading, bullets in sections:
        flat.append(("HEAD", heading))
        for b in bullets:
            if b:
                flat.append((None, b))
    if not flat:
        return

    idx = 0
    page = 0
    while idx < len(flat):
        slide = _blank(prs)
        _title_bar(slide, title if page == 0 else f"{title} (продолжение)")
        tf = _body(slide)
        used = 0
        count = 0
        first = True
        while idx < len(flat) and count < _MAX_BULLETS and used < _CHAR_BUDGET:
            kind, text = flat[idx]
            if kind == "HEAD":
                # Avoid a heading stranded at the very bottom of a slide.
                if not first and used > _CHAR_BUDGET - 200:
                    break
                _para(tf, text, size=14, bold=True, color=_ACCENT, first=first)
            else:
                _para(tf, text, bullet=True, first=first)
                count += 1
            first = False
            used += len(text) + 30
            idx += 1
        page += 1


def _style_cell(cell, *, text: str = "", bold: bool = False, size: int = 11, color=_DARK, fill=None, align=PP_ALIGN.LEFT) -> None:
    if text:
        cell.text = text
    tf = cell.text_frame
    tf.word_wrap = True
    for paragraph in tf.paragraphs:
        paragraph.alignment = align
        for run in paragraph.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill


def _fmt_level(value: Any) -> str:
    if value in (None, ""):
        return "—"
    try:
        num = float(value)
    except (TypeError, ValueError):
        return str(value)
    text = f"{num:.1f}".rstrip("0").rstrip(".")
    return text.replace(".", ",")
